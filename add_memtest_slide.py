"""
슬라이드 추가:
  - 슬라이드 20: Microwire 읽기 수정 + 메모리 테스트(State C)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPT_PATH = r"D:\3_STM32CubeIDE\93c56\93c56_overview.pptx"

C_TITLE  = RGBColor(0x1F, 0x49, 0x7D)
C_HEAD   = RGBColor(0x2E, 0x74, 0xB5)
C_OK     = RGBColor(0x37, 0x86, 0x38)
C_FAIL   = RGBColor(0xC0, 0x00, 0x00)
C_WARN   = RGBColor(0xED, 0x7D, 0x31)
C_GRAY   = RGBColor(0x40, 0x40, 0x40)
C_BG     = RGBColor(0xF2, 0xF2, 0xF2)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF8, 0xF8, 0xF8)
C_BLUE   = RGBColor(0x00, 0x70, 0xC0)

def tb(slide, text, left, top, width, height,
       sz=12, bold=False, color=C_GRAY,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return box

def rect(slide, left, top, width, height, fill, line=None):
    s = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def slide_title(slide, title, subtitle=""):
    rect(slide, 0, 0, 10, 0.7, C_TITLE)
    tb(slide, title, 0.15, 0.05, 9.0, 0.6,
       sz=20, bold=True, color=C_WHITE)
    if subtitle:
        tb(slide, subtitle, 0.15, 0.72, 9.5, 0.35,
           sz=12, color=C_HEAD, italic=True)

def section_box(slide, label, left, top, width, color=C_HEAD):
    rect(slide, left, top, width, 0.30, color)
    tb(slide, label, left+0.08, top+0.04, width-0.1, 0.25,
       sz=11, bold=True, color=C_WHITE)

# ── PPT 열기 ──────────────────────────────────────────────────
prs = Presentation(PPT_PATH)
layout = prs.slide_layouts[0]
n = len(prs.slides)
print(f"기존 슬라이드 수: {n}")

# ══════════════════════════════════════════════════════════════
# 슬라이드 N+1: Microwire 읽기 수정 + 메모리 테스트
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(layout)
slide_title(s,
    f"슬라이드 {n+1}  Microwire 읽기 수정 + 메모리 테스트(State C)",
    "2026.04 | 0xAA/0x55 R/W 정상 동작 확인 · 전체 주소 검증 단계 추가")

# ── 왼쪽: Microwire 읽기 버그 수정 ──────────────────────────
section_box(s, "① Microwire 읽기 버그 수정 (EE_Read)", 0.3, 0.85, 4.6, C_FAIL)

bug_items = [
    ("문제",  "EE_Read()에서 ee_recv_bit() 로 더미비트 스킵\n"
              "→ D7 을 소비해 데이터 1비트 좌이동 발생"),
    ("증상",  "Write 0x5A  →  Read 0xB4  (0x5A << 1)\n"
              "Write 0xAA  →  Read 0x4A  등 항상 불일치"),
    ("수정",  "ee_recv_bit() 제거\n"
              "→ ee_delay() 로 교체 (tV ≤ 200ns 확보)"),
    ("결과",  "0xAA / 0x55 정상 읽기 확인\n"
              "Write == Read 일치"),
]
y = 1.20
for label, desc in bug_items:
    rect(s, 0.3, y, 1.1, 0.22, C_HEAD)
    tb(s, label, 0.35, y+0.03, 1.0, 0.18, sz=9, bold=True, color=C_WHITE)
    rect(s, 1.40, y, 3.50, 0.22, C_LGRAY, C_HEAD)
    tb(s, desc, 1.48, y+0.02, 3.36, 0.22, sz=8, color=C_GRAY)
    y += 0.27

# 코드 비교 박스
y += 0.05
rect(s, 0.3, y, 4.6, 0.22, RGBColor(0xC0,0x00,0x00))
tb(s, "수정 전 (버그)", 0.35, y+0.03, 4.4, 0.18, sz=9, bold=True, color=C_WHITE)
y += 0.22
rect(s, 0.3, y, 4.6, 0.28, RGBColor(0xFF,0xEB,0xEB), C_FAIL)
tb(s, "ee_recv_bit();   /* Dummy bit — D7 소비! */",
   0.38, y+0.04, 4.40, 0.22, sz=8.5, color=C_FAIL)
y += 0.33

rect(s, 0.3, y, 4.6, 0.22, C_OK)
tb(s, "수정 후 (정상)", 0.35, y+0.03, 4.4, 0.18, sz=9, bold=True, color=C_WHITE)
y += 0.22
rect(s, 0.3, y, 4.6, 0.28, RGBColor(0xE8,0xF5,0xE8), C_OK)
tb(s, "ee_delay();      /* D7 드라이브 대기 (tV ≤ 200ns) */",
   0.38, y+0.04, 4.40, 0.22, sz=8.5, color=C_OK)

# ── 오른쪽: 메모리 테스트 State C ───────────────────────────
section_box(s, "② 메모리 테스트 State C 추가", 5.2, 0.85, 4.6, C_OK)

flow = [
    (C_BLUE, "State B 완료 후 자동 실행"),
    (C_OK,   "Pass 1: addr 0~255 에 0xAA 쓰기"),
    (C_OK,   "         addr 0~255 에서 0xAA 읽기 검증"),
    (C_OK,   "Pass 2: addr 0~255 에 0x55 쓰기"),
    (C_OK,   "         addr 0~255 에서 0x55 읽기 검증"),
    (C_WARN, "실패 시 → 주소 / 읽기값 UART 출력"),
    (C_FAIL, "FAIL 칩 → State D 프로그래밍 SKIP"),
]
y2 = 1.20
for clr, txt in flow:
    rect(s, 5.2, y2, 0.08, 0.20, clr)
    tb(s, txt, 5.35, y2+0.02, 4.35, 0.20, sz=9, color=C_GRAY)
    y2 += 0.24

# UART 출력 예시
y2 += 0.05
rect(s, 5.2, y2, 4.6, 0.22, C_TITLE)
tb(s, "UART 출력 예시", 5.28, y2+0.03, 4.4, 0.18, sz=9, bold=True, color=C_WHITE)
y2 += 0.22

uart_ok = [
    ("[TEST] === Memory Test (0xAA / 0x55) ===", C_HEAD),
    ("[TEST] CS00: PASS",                        C_OK),
    ("[TEST] CS01: PASS",                        C_OK),
    ("[TEST] CS03: 0xAA RD ERR  addr=47  got=0xB4", C_FAIL),
    ("[TEST] CS03: FAIL",                        C_FAIL),
    ("[TEST] Result: PASS=11  FAIL=1",           C_WARN),
    ("[TEST] Failed chips: CS03",                C_WARN),
]
rect(s, 5.2, y2, 4.6, 0.01 + len(uart_ok)*0.24, C_BG, C_HEAD)
for line, clr in uart_ok:
    tb(s, line, 5.30, y2+0.02, 4.42, 0.22, sz=8, color=clr)
    y2 += 0.24

# ── 하단: LED 동작 요약 ────────────────────────────────────
rect(s, 0.3, 6.50, 9.4, 0.28, C_TITLE)
tb(s,
   "테스트 중 → ■ 파란색(BLUE)    PASS → ■ 녹색(GREEN)    FAIL → ■ 적색(RED) + 프로그래밍 SKIP",
   0.45, 6.53, 9.1, 0.24, sz=10, bold=True, color=C_WHITE)

rect(s, 0.3, 6.82, 9.4, 0.28, C_OK)
tb(s,
   "최종 결과:  메모리 R/W 정상 동작 확인  ·  0xAA / 0x55 전체 256주소 검증  ·  불량 칩 자동 분류",
   0.45, 6.85, 9.1, 0.24, sz=10, bold=True, color=C_WHITE)

prs.save(PPT_PATH)
print(f"PPT 저장 완료: {PPT_PATH}")
print(f"총 슬라이드: {len(prs.slides)}")

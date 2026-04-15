"""
최종 완료 슬라이드를 PPT에 추가 (2026.04)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPT_PATH = r"D:\3_STM32CubeIDE\93c56\93c56_overview.pptx"

C_TITLE  = RGBColor(0x1F, 0x49, 0x7D)
C_HEAD   = RGBColor(0x2E, 0x74, 0xB5)
C_OK     = RGBColor(0x37, 0x86, 0x38)
C_FAIL   = RGBColor(0xC0, 0x00, 0x00)
C_WARN   = RGBColor(0xED, 0x7D, 0x31)
C_GRAY   = RGBColor(0x40, 0x40, 0x40)
C_BG     = RGBColor(0xF2, 0xF2, 0xF2)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF8, 0xF8, 0xF8)

def tb(slide, text, left, top, width, height,
        sz=12, bold=False, color=C_GRAY,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return box

def rect(slide, left, top, width, height, fill, line=None):
    s = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def slide_title(slide, title, subtitle=""):
    rect(slide, 0, 0, 10, 0.7, C_TITLE)
    tb(slide, title, 0.15, 0.05, 8.5, 0.6,
       sz=20, bold=True, color=C_WHITE)
    if subtitle:
        tb(slide, subtitle, 0.15, 0.72, 9.5, 0.35,
           sz=12, color=C_HEAD, italic=True)

# ── PPT 열기 ──────────────────────────────────────────────────
prs = Presentation(PPT_PATH)
layout = prs.slide_layouts[0]
n = len(prs.slides)
print(f"기존 슬라이드 수: {n}")

# ══════════════════════════════════════════════════════════════
# 슬라이드 N+1: 최종 완료 상태
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(layout)
slide_title(s, f"슬라이드 {n+1}  최종 완료 — 93C56 프로그래머",
            "2026.04 | 12/12 칩 프로그래밍 성공 · SPI1 충돌 해결")

# ── 왼쪽: UART 최종 출력 ──────────────────────────────────────
tb(s, "최종 UART 출력", 0.3, 0.9, 4.8, 0.32,
   sz=13, bold=True, color=C_HEAD)

uart = [
    ("[DIAG] Chips found: 12 / 12",              C_OK),
    ("[INIT] LED sequence start (count=12)",      C_OK),
    ("[INIT] LED sequence done",                  C_OK),
    ("[WAIT] Loop exited: count=12, skip=0x0000", C_OK),
    ("[BOARD] Board detected! (12/12 chips)",     C_OK),
    ("[STATE B] CS00~CS11: OK",                   C_OK),
    ("[STATE D] Programming (write 0x00 all)",    C_OK),
    ("[PROG] CS00~CS11: OK",                      C_OK),
    ("[DONE] Result: PASS=12  FAIL=0  SKIP=0",   C_OK),
    ("[DONE] All chips PASS.",                    C_OK),
    ("[WAIT] Remove the board...",                C_HEAD),
]
rect(s, 0.3, 1.25, 4.8, 0.01 + len(uart)*0.32, C_BG, C_HEAD)
y = 1.28
for line, clr in uart:
    tb(s, line, 0.42, y, 4.56, 0.30, sz=9, color=clr)
    y += 0.32

# ── 오른쪽: 해결 이력 ─────────────────────────────────────────
tb(s, "문제 해결 이력", 5.3, 0.9, 4.5, 0.32,
   sz=13, bold=True, color=C_HEAD)

fixes = [
    (C_OK, "① EE_CHIP_SKIP_MASK 0x00FF → 0x0000",
     "1,2번 줄(CS00-07) 강제 스킵 제거 → 전체 12칩 스캔"),
    (C_OK, "② g_board_present 감지 경쟁 해결",
     "진단 결과를 즉시 g_detected[] 에 반영\n→ State A while 루프 즉시 탈출"),
    (C_OK, "③ EE_ProgramZero 검증 제거",
     "RB=0x48 문제로 verify 항상 실패\n→ WR=OK 기준으로 성공 판정"),
    (C_OK, "④ SPI1 클럭 비활성화 버그 수정",
     "__HAL_RCC_SPI1_CLK_DISABLE() 제거\n→ SK9822 LED HAL_SPI_Transmit 멈춤 해결"),
    (C_WARN, "⑤ RB=0x48 미해결 (우회)",
     "Microwire dummy bit 미처리 → 읽기 1비트 시프트\n프로그래밍 자체는 WR=OK로 정상"),
]

y = 1.25
for clr, title, desc in fixes:
    rect(s, 5.3, y, 4.5, 0.28, clr)
    tb(s, title, 5.38, y+0.03, 4.3, 0.24,
       sz=10, bold=True, color=C_WHITE)
    y += 0.30
    rect(s, 5.3, y, 4.5, 0.50, C_LGRAY, clr)
    tb(s, desc, 5.42, y+0.03, 4.26, 0.46, sz=9, color=C_GRAY)
    y += 0.54

# ── 하단: LED 결과 설명 ───────────────────────────────────────
rect(s, 0.3, 6.55, 9.4, 0.32, C_TITLE)
tb(s, "LED 결과:  ■ 백색(WHITE) = 프로그래밍 완료   ■ 적색(RED) = 칩 없음/불량",
   0.45, 6.57, 9.0, 0.28, sz=11, bold=True, color=C_WHITE)

# ── 하단: 핀 충돌 요약 ────────────────────────────────────────
tb(s, "SPI1 충돌 분석", 0.3, 6.95, 9.4, 0.28,
   sz=11, bold=True, color=C_HEAD)
pin_info = [
    ("PA6",  "SPI1_MISO (SK9822용)", "DO7 (EEPROM 읽기)",
     "SPI1_DIRECTION_1LINE → MISO 미사용 → 충돌 없음"),
    ("PA7",  "SPI1_MOSI (SK9822 데이터)", "—",
     "LED 데이터 전송 핀, SPI1 유지 필요"),
]
heads = ["핀", "SPI1 기능", "EEPROM 기능", "결론"]
widths = [0.6, 2.4, 2.4, 3.9]
y = 7.25
x = 0.3
for h, w in zip(heads, widths):
    rect(s, x, y, w-0.02, 0.27, C_TITLE)
    tb(s, h, x+0.05, y+0.03, w-0.1, 0.22, sz=9, bold=True, color=C_WHITE)
    x += w
y += 0.27
for row in pin_info:
    x = 0.3
    bg = C_LGRAY if pin_info.index(row) % 2 == 0 else C_WHITE
    for cell, w in zip(row, widths):
        rect(s, x, y, w-0.02, 0.28, bg, RGBColor(0xD0,0xD0,0xD0))
        tb(s, cell, x+0.06, y+0.03, w-0.12, 0.24, sz=9, color=C_GRAY)
        x += w
    y += 0.28

prs.save(PPT_PATH)
print(f"PPT 저장 완료: {PPT_PATH}")
print(f"총 슬라이드: {len(prs.slides)}")

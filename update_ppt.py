"""
update_ppt.py  —  B/W, 배경 없음 버전  (7클립 최종)
1) 93c56_overview.pptx  : Slide 12 음성 업데이트 + Slide 21 Phase2 업데이트
2) 93c56_user_manual.pptx : 사용자 매뉴얼 재생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.33)
H = Inches(7.5)

BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY1  = RGBColor(0xF2, 0xF2, 0xF2)
GRAY2  = RGBColor(0xD9, 0xD9, 0xD9)
DGRAY  = RGBColor(0x40, 0x40, 0x40)
LINE   = RGBColor(0xBB, 0xBB, 0xBB)


def clear_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def tb(slide, text, l, t, w, h,
       size=14, bold=False, italic=False,
       color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def hline(slide, top, left=Inches(0.5), width=Inches(12.3), thick=Pt(1)):
    ln = slide.shapes.add_shape(1, left, top, width, Inches(0.02))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE
    ln.line.fill.background()


def add_title(slide, title, subtitle=None):
    tb(slide, title,
       Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.7),
       size=28, bold=True, color=BLACK)
    hline(slide, Inches(0.98))
    if subtitle:
        tb(slide, subtitle,
           Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.38),
           size=13, italic=True, color=DGRAY)


def add_table(slide, headers, rows, left, top, width, height, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = cw

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = GRAY2
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = BLACK

    for ri, row in enumerate(rows):
        bg = GRAY1 if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.size = Pt(12); run.font.color.rgb = BLACK
    return tbl


def clear_slide(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


# ═══════════════════════════════════════════════════════════════
# Part 1 — 기존 PPT 업데이트
# ═══════════════════════════════════════════════════════════════
prs = Presentation('D:/3_STM32CubeIDE/93c56/93c56_overview.pptx')

# ── Slide 12: 음성 메시지 안내 (7종, 16kHz) ──────────────────
slide12 = prs.slides[11]
clear_slide(slide12)
clear_bg(slide12)

add_title(slide12, "음성 메시지 안내 (7종)",
          "gen_voice.py — Windows TTS (Heami KO)  ·  16 kHz · 16-bit mono PCM")

headers = ["#", "음성 내용", "사용 시점"]
rows = [
    ("0",   "디아지용 팁 보드를 올려주세요",  "상태 A: LED 쇼 중 최초 안내"),
    ("1",   "프로그램이 완료되었습니다",       "프로그래밍 완료 후"),
    ("2",   "적색 표시는 불량입니다",          "불량 칩 존재 시"),
    ("3",   "보드를 분리하여 주십시요",        "상태 E: 보드 제거 요청 (30초 반복)"),
    ("4",   "보드가 분리되었습니다",           "상태 E: 보드 제거 감지 직후"),
    ("5 *", "새 보드를 올려 주십시요",         "보드 분리 후 + 10초마다 반복"),
    ("6 *", "프로그램을 시작합니다",           "2초 안정화 대기 후 프로그래밍 직전"),
]
add_table(slide12, headers, rows,
          Inches(0.5), Inches(1.55), Inches(12.3), Inches(5.0),
          [Inches(0.65), Inches(7.0), Inches(4.65)])

tb(slide12, "* = Phase 2 신규 추가  /  8 kHz -> 16 kHz 업그레이드 (자음 명료도 향상)  /  INTRO·IS_93C56·IS_93C46 3종 제거",
   Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.35),
   size=11, italic=True, color=DGRAY)

# ── Slide 21: Phase 2 주요 업데이트 ──────────────────────────
# 이미 존재하면 내용만 교체, 없으면 신규 추가
if len(prs.slides) > 20:
    s21 = prs.slides[20]
    clear_slide(s21)
else:
    s21 = prs.slides.add_slide(prs.slide_layouts[0])

clear_bg(s21)
add_title(s21, "Phase 2 주요 업데이트 (2026.04)",
          "대량 프로그래밍 모드 안정화  ·  음성 개선  ·  보드 감지 강화")

sections = [
    [
        ("1. 보드 감지 개선", [
            ("HAL_GetTick() -> osKernelGetTickCount()",
             "FreeRTOS 환경에서 HAL 타이머가 0 반환 -> 쓰로틀 무한 실패 수정"),
            ("EE_Write 타임아웃 기반 존재 감지",
             "이미 프로그래밍된 보드 재삽입 시 데이터 비교 오판 방지"),
        ]),
        ("2. 2초 안정화 대기", [
            ("보드 감지 후 즉시 프로그래밍 금지",
             "파란 LED 전체 2초 점멸 -> 사용자 보드 가압으로 접촉 확보"),
            ("2초 후 연결 재확인 (board_detect)",
             "재확인 실패 시 State A 자동 재시작 (continue)"),
        ]),
    ],
    [
        ("3. 음성 최적화 (최종 7종)", [
            ("8 kHz -> 16 kHz 업그레이드",
             "자음 (ㅅ/ㅈ/ㅊ) 명료도 향상 / Flash 사용량 조정"),
            ("불필요 3종 제거 · 신규 2종 추가",
             "제거: INTRO·IS_93C56·IS_93C46 / 추가: #5 새 보드 · #6 프로그램 시작"),
        ]),
        ("4. 대기 화면 반복 안내", [
            ("보드 미감지 상태 지속 시",
             "10초마다 '새 보드를 올려 주십시요' 자동 반복 출력"),
            ("8가지 LED 쇼 패턴 내 체크포인트",
             "각 패턴 종료 시점마다 10초 경과 여부 판정"),
        ]),
    ],
]

for col_idx, col_sections in enumerate(sections):
    lft = Inches(0.5) + col_idx * Inches(6.5)
    y   = Inches(1.55)
    for sec_title, items in col_sections:
        tb(s21, sec_title, lft, y, Inches(6.1), Inches(0.38),
           size=14, bold=True, color=BLACK)
        y += Inches(0.4)
        hline(s21, y, left=lft, width=Inches(6.0))
        y += Inches(0.1)
        for item_title, item_desc in items:
            tb(s21, "- " + item_title, lft+Inches(0.1), y, Inches(5.9), Inches(0.32),
               size=12, bold=True, color=BLACK)
            y += Inches(0.3)
            tb(s21, "  " + item_desc, lft+Inches(0.1), y, Inches(5.9), Inches(0.38),
               size=11, color=DGRAY)
            y += Inches(0.42)
        y += Inches(0.2)

vl = s21.shapes.add_shape(1, Inches(6.6), Inches(1.55), Inches(0.02), Inches(5.6))
vl.fill.solid(); vl.fill.fore_color.rgb = LINE
vl.line.fill.background()

prs.save('D:/3_STM32CubeIDE/93c56/93c56_overview.pptx')
print("93c56_overview.pptx 저장 완료")


# ═══════════════════════════════════════════════════════════════
# Part 2 — 사용자 매뉴얼 PPT 재생성  (B/W, 배경 없음, 7클립)
# ═══════════════════════════════════════════════════════════════
man = Presentation()
man.slide_width  = W
man.slide_height = H
blank = man.slide_layouts[6]


def ms(layout=None):
    s = man.slides.add_slide(layout or blank)
    clear_bg(s)
    return s


# ── 슬라이드 1: 표지 ──────────────────────────────────────────
s = ms()
for y_pos, thick in [(Inches(1.9), Inches(0.04)), (Inches(5.5), Inches(0.04))]:
    ln = s.shapes.add_shape(1, Inches(1.5), y_pos, Inches(10.3), thick)
    ln.fill.solid(); ln.fill.fore_color.rgb = BLACK
    ln.line.fill.background()

tb(s, "DAGI Tip Programmer",
   Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.1),
   size=40, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
tb(s, "93C46 EEPROM 대량 프로그래밍 시스템",
   Inches(1.5), Inches(3.15), Inches(10.3), Inches(0.6),
   size=20, color=BLACK, align=PP_ALIGN.CENTER)
tb(s, "사 용 설 명 서",
   Inches(1.5), Inches(3.85), Inches(10.3), Inches(0.55),
   size=18, bold=True, color=DGRAY, align=PP_ALIGN.CENTER)
tb(s, "DSI INC.",
   Inches(1.5), Inches(6.7), Inches(10.3), Inches(0.4),
   size=14, color=DGRAY, align=PP_ALIGN.CENTER)

# ── 슬라이드 2: 제품 개요 ──────────────────────────────────────
s = ms()
add_title(s, "제품 개요")

specs = [
    ("제품명",   "DAGI Tip Programmer (93C46 대량 프로그래밍 시스템)"),
    ("목적",     "DAGI Tip PCB 보드의 93C46 EEPROM 12개 동시 프로그래밍"),
    ("제어 MCU", "STM32H562RGT6 @ 32 MHz  /  FreeRTOS"),
    ("지원 칩",  "93C46  (1 Kbit · 128x8 bit · 7-bit 주소)"),
    ("슬롯 수",  "12슬롯 (CS00 ~ CS11) 동시 프로그래밍"),
    ("음성 안내","한국어 TTS 7종  /  16 kHz · 16-bit mono PCM"),
    ("LED 표시", "SK9822 RGB LED 12개 — 슬롯별 진행 / 결과 표시"),
    ("디버그",   "UART1 115200 bps  (PB14 TX / PB15 RX)"),
]
y = Inches(1.55)
for i, (k, v) in enumerate(specs):
    bg = GRAY1 if i % 2 == 0 else WHITE
    row_box = s.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Inches(0.43))
    row_box.fill.solid(); row_box.fill.fore_color.rgb = bg
    row_box.line.fill.background()
    tb(s, k, Inches(0.6), y+Inches(0.05), Inches(2.3), Inches(0.36),
       size=13, bold=True, color=BLACK)
    tb(s, v, Inches(3.0), y+Inches(0.05), Inches(9.6), Inches(0.36),
       size=13, color=BLACK)
    y += Inches(0.45)

# ── 슬라이드 3: 사용 전 준비 ──────────────────────────────────
s = ms()
add_title(s, "사용 전 준비")

items = [
    ("①", "전원 연결",
     "프로그래머 보드에 DC 전원을 공급합니다."),
    ("②", "빈 보드 준비",
     "EEPROM이 장착된 DAGI Tip PCB 보드를 준비합니다."),
    ("③", "LED 패턴 확인",
     "12개 LED가 순차 점등 쇼를 반복하며 보드 삽입 대기 상태임을 알립니다."),
    ("④", "음성 안내",
     "'디아지용 팁 보드를 올려주세요' 후 10초마다 '새 보드를 올려 주십시요'가 반복됩니다."),
]
y = Inches(1.55)
for num, title, desc in items:
    circ = s.shapes.add_shape(9, Inches(0.5), y+Inches(0.05), Inches(0.42), Inches(0.42))
    circ.fill.solid(); circ.fill.fore_color.rgb = GRAY2
    circ.line.color.rgb = BLACK; circ.line.width = Pt(0.75)
    tb(s, num, Inches(0.5), y+Inches(0.06), Inches(0.42), Inches(0.38),
       size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    tb(s, title, Inches(1.1), y, Inches(11.5), Inches(0.36),
       size=14, bold=True, color=BLACK)
    tb(s, desc,  Inches(1.1), y+Inches(0.36), Inches(11.5), Inches(0.38),
       size=12, color=DGRAY)
    y += Inches(0.9)

# ── 슬라이드 4: 프로그래밍 절차 ───────────────────────────────
s = ms()
add_title(s, "프로그래밍 절차 (단계별)")

steps = [
    ("STEP 1", "보드 삽입",
     "DAGI Tip 보드를 프로그래머 소켓에 올려 놓습니다.\n보드 감지 시 파란 LED 전체 점멸 시작."),
    ("STEP 2", "2초 안정화 대기",
     "파란 LED 점멸 약 2초 동안 보드를 가볍게\n눌러 접촉을 확보해 주세요."),
    ("STEP 3", "연결 재확인 (자동)",
     "2초 후 12개 EEPROM 접촉 자동 확인.\n접촉 불량 시 STEP 1로 자동 재시작."),
    ("STEP 4", "프로그래밍 시작",
     "'프로그램을 시작합니다' 음성 후\nCS00~CS11 순서로 기록 진행."),
    ("STEP 5", "결과 확인",
     "흰색 LED = 정상 완료\n적색 LED = 불량 / 미접촉"),
    ("STEP 6", "보드 분리 후 교체",
     "'보드를 분리하여 주십시요' 안내 후 제거.\n'새 보드를 올려 주십시요' 안내 -> STEP 1 반복."),
]
cols = 3
for idx, (step, title, desc) in enumerate(steps):
    col = idx % cols
    row = idx // cols
    lft = Inches(0.4)  + col * Inches(4.35)
    tp  = Inches(1.55) + row * Inches(2.8)
    box = s.shapes.add_shape(1, lft, tp, Inches(4.1), Inches(2.55))
    box.fill.solid(); box.fill.fore_color.rgb = GRAY1
    box.line.color.rgb = BLACK; box.line.width = Pt(0.75)
    tb(s, step,  lft+Inches(0.1), tp+Inches(0.08), Inches(1.5), Inches(0.3),
       size=11, bold=True, color=DGRAY)
    tb(s, title, lft+Inches(0.1), tp+Inches(0.38), Inches(3.9), Inches(0.38),
       size=15, bold=True, color=BLACK)
    tb(s, desc,  lft+Inches(0.1), tp+Inches(0.82), Inches(3.9), Inches(1.5),
       size=12, color=DGRAY)

# ── 슬라이드 5: LED 표시 의미 ──────────────────────────────────
s = ms()
add_title(s, "LED 표시 의미")

led_h = ["LED 색상", "상태", "설명"]
led_r = [
    ("파란색 점멸", "보드 감지 / 안정화 대기",   "보드 삽입 감지 후 2초 안정화 중"),
    ("파란색 점등", "프로그래밍 진행 중",         "해당 슬롯 EEPROM 기록 중"),
    ("흰색 점등",   "프로그래밍 완료 (정상)",     "해당 슬롯 Write + Verify 성공"),
    ("적색 점등",   "불량 / 미접촉",              "Write 실패 또는 칩 없음"),
    ("무지개 쇼",   "대기 중 (보드 없음)",        "8가지 순환 패턴 — 보드 삽입 대기 상태"),
    ("전체 소등",   "초기화 중 / 전원 OFF",       "—"),
]
add_table(s, led_h, led_r,
          Inches(0.5), Inches(1.55), Inches(12.3), Inches(4.5),
          [Inches(2.5), Inches(3.8), Inches(6.0)])

tb(s, "※ 각 슬롯(CS00~CS11)의 LED가 독립적으로 상태를 표시합니다.",
   Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
   size=12, italic=True, color=DGRAY)

# ── 슬라이드 6: 음성 안내 목록 (7종) ──────────────────────────
s = ms()
add_title(s, "음성 안내 목록 (7종)",
          "16 kHz · 16-bit mono PCM  /  한국어 TTS (Microsoft Heami)")

v_h = ["#", "음성 내용", "출력 시점"]
v_r = [
    ("0",  "디아지용 팁 보드를 올려주세요",  "시작 시 1회"),
    ("1",  "프로그램이 완료되었습니다",       "프로그래밍 완료"),
    ("2",  "적색 표시는 불량입니다",          "불량 칩 존재 시"),
    ("3",  "보드를 분리하여 주십시요",        "완료 후 보드 제거 요청 (30초 반복)"),
    ("4",  "보드가 분리되었습니다",           "보드 제거 감지"),
    ("5",  "새 보드를 올려 주십시요",         "보드 분리 후 + 10초마다 반복"),
    ("6",  "프로그램을 시작합니다",           "2초 안정화 후 프로그래밍 직전"),
]
add_table(s, v_h, v_r,
          Inches(0.5), Inches(1.55), Inches(12.3), Inches(4.8),
          [Inches(0.55), Inches(7.3), Inches(4.45)])

# ── 슬라이드 7: 주의사항 ───────────────────────────────────────
s = ms()
add_title(s, "주의사항 및 오류 대처")

warnings = [
    ("[ 주의 1 ]  보드 삽입 시 접촉 확인",
     "보드 삽입 후 파란 LED 점멸 2초 동안 보드를 가볍게 눌러 접촉을 확보하세요.\n"
     "접촉 불량 -> 연결 재확인 실패 -> 자동으로 대기 화면(State A)으로 복귀합니다."),
    ("[ 주의 2 ]  적색 LED 슬롯 확인",
     "프로그래밍 완료 후 적색 LED 슬롯은 EEPROM 불량 또는 미접촉 상태입니다.\n"
     "해당 위치의 EEPROM 및 소켓 접촉 상태를 확인하세요."),
    ("[ 주의 3 ]  프로그래밍 중 전원 차단 금지",
     "프로그래밍 진행 중(파란 LED 점등) 전원을 끊으면 EEPROM 데이터가 손상됩니다."),
    ("[ 안내 ]  연속 프로그래밍",
     "보드 분리 -> '새 보드를 올려 주십시요' 안내 -> 다음 보드 삽입 -> STEP 1 반복.\n"
     "전원을 끄지 않고 보드 교체만으로 연속 작업이 가능합니다."),
]
y = Inches(1.55)
for title, desc in warnings:
    box = s.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Inches(1.3))
    box.fill.solid(); box.fill.fore_color.rgb = GRAY1
    box.line.color.rgb = LINE; box.line.width = Pt(0.5)
    tb(s, title, Inches(0.65), y+Inches(0.08), Inches(11.9), Inches(0.36),
       size=13, bold=True, color=BLACK)
    tb(s, desc,  Inches(0.65), y+Inches(0.44), Inches(11.9), Inches(0.72),
       size=12, color=DGRAY)
    y += Inches(1.38)

man.save('D:/3_STM32CubeIDE/93c56/93c56_user_manual.pptx')
print("93c56_user_manual.pptx 저장 완료")
print("완료!")
